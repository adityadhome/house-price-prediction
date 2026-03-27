import os
import sys

# Try importing the required libraries, install them if they fail
try:
    from docx import Document
    from docx.shared import Pt, Inches
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
except ImportError:
    os.system(f"{sys.executable} -m pip install python-docx python-pptx")
    from docx import Document
    from docx.shared import Pt, Inches
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt

# 1. Generate MS Word Document (.docx)
def generate_word():
    doc = Document()
    doc.add_heading('PROJECT DOCUMENTATION', 0)
    doc.add_heading('Pune Real Estate Analytics & Prediction', level=1)

    doc.add_heading('1. ABSTRACT', level=2)
    doc.add_paragraph(
        "The Real Estate industry is one of the most unpredictable and highly volatile sectors. "
        "Predicting property prices manually is prone to human error and bias. The primary objective of this project "
        "is to analyze historical real estate data and develop a Machine Learning (ML) model capable of predicting "
        "house prices with high accuracy. The project implements an interactive Web Dashboard predicting prices by selecting "
        "specific areas (e.g., Pune areas) and features seamlessly."
    )

    doc.add_heading('2. INTRODUCTION', level=2)
    doc.add_paragraph(
        "With globalization and rapid urbanization, investing in real estate has become highly lucrative. "
        "However, buyers and sellers often face a pricing dilemma. This project leverages Multiple Linear Regression "
        "so the system learns the impact of housing attributes (bedrooms, bathrooms, sqft). A clean Front-End is provided "
        "via the Streamlit framework."
    )

    doc.add_heading('3. SYSTEM REQUIREMENTS', level=2)
    doc.add_heading('Hardware Requirements', level=3)
    doc.add_paragraph("• Processor: Minimum Intel Core i3 / AMD Ryzen 3 or above.\n• RAM: 4 GB (8 GB recommended)\n• Storage: 1 GB free space.")
    doc.add_heading('Software Requirements', level=3)
    doc.add_paragraph("• OS: Windows 10/11\n• Language: Python 3.8+\n• Libraries: pandas, numpy, scikit-learn, matplotlib, streamlit.")

    doc.add_heading('4. SYSTEM ARCHITECTURE & METHODOLOGY', level=2)
    doc.add_paragraph("1. Data Processing & Cleaning\n2. Feature Engineering (One-Hot Encoding for Pune Cities)\n3. Training the Model (80-20 Split)\n4. The Algorithm (Linear Regression)")

    doc.add_heading('5. IMPLEMENTATION DETAILS', level=2)
    doc.add_paragraph(
        "The dashboard is built via Streamlit. A select box allows localized City/Area selection (Baner, Viman Nagar). "
        "Missing Dummy columns are matched via reindex(), and the pred() function executes the estimation. Output is in USD and INR."
    )

    doc.add_heading('6. RESULTS & EVALUATION', level=2)
    doc.add_paragraph(
        "The model produces impressive results during evaluation against the 20% test data. "
        "We used MAE, RMSE, and R-squared to measure performance accuracy. The plot 'prediction_accuracy.png' was generated successfully."
    )

    doc.add_heading('7. CONCLUSION', level=2)
    doc.add_paragraph(
        "This project successfully converges Machine Learning with Web Engineering. "
        "The application runs robustly and can instantly formulate property estimations removing any requirement for manual valuation."
    )

    doc.save('Pune_House_Price_Prediction_Report.docx')
    print("Word Document Generated Successfully!")

# 2. Generate PowerPoint Presentation (.pptx)
def generate_ppt():
    prs = Presentation()

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Pune Real Estate Analytics"
    slide.placeholders[1].text = "House Price Prediction using Machine Learning\nClass Project Submission"

    # Slide 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction"
    tf = slide.placeholders[1].text_frame
    tf.text = "• The real estate market is highly volatile."
    p = tf.add_paragraph()
    p.text = "• Our Solution: An AI & ML based system predicting house prices in Pune."
    p = tf.add_paragraph()
    p.text = "• Goal: Help users estimate the fair market value instantly."

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Statement"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Buyers overpay and sellers underprice due to lack of data."
    p = tf.add_paragraph()
    p.text = "• Manual price estimation is tedious and inaccurate."
    p = tf.add_paragraph()
    p.text = "• Need: A fast, automated system returning exact prices based on features."

    # Slide 4: Objectives & Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Objectives & Tech Stack"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Objective: Process historical data, train an ML model, and build a Web App."
    p = tf.add_paragraph()
    p.text = "• Language: Python"
    p = tf.add_paragraph()
    p.text = "• Libraries: Scikit-Learn, Pandas, NumPy, Matplotlib, Streamlit"

    # Slide 5: Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology & Workflow"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Data Collection & Cleaning (Dropping null values)"
    p = tf.add_paragraph()
    p.text = "2. Feature Engineering (One-Hot Encoding for Pune areas)"
    p = tf.add_paragraph()
    p.text = "3. Train-Test Split (80% / 20%)"
    p = tf.add_paragraph()
    p.text = "4. Linear Regression Model Training"

    # Slide 6: Results & Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Results & Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Model accurately maps housing features to exact market prices."
    p = tf.add_paragraph()
    p.text = "• The interactive Web Dashboard displays inputs cleanly and outputs ₹ INR."
    p = tf.add_paragraph()
    p.text = "• Conclusion: A robust integration of Data Science & Web Engineering!"

    # Slide 7: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Thank You!"
    slide.placeholders[1].text = "Questions?"

    prs.save('Pune_Real_Estate_Presentation.pptx')
    print("PowerPoint Presentation Generated Successfully!")

if __name__ == "__main__":
    generate_word()
    generate_ppt()
